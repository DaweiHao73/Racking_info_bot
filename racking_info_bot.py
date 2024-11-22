# 導入必要的庫
import pandas as pd  # 用於處理Excel數據
from pptx import Presentation  # 用於創建PowerPoint
from pptx.util import Inches, Pt  # 用於設置PowerPoint中的尺寸和字體大小
from pptx.enum.text import PP_ALIGN  # 用於設置文字對齊方式
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT  # 新增：用於段落對齊

try:
    # 設定Excel文件路徑並讀取數據
    file_path = r"C:\Users\dwhao\OneDrive - The Dairy Farm Company Ltd\HCS VM Stuff\IMS\各店貨架借用表.xlsx"
    # 讀取指定工作表的數據到DataFrame
    df = pd.read_excel(file_path, sheet_name="借出 & 借用表")
    
    # 創建新的PowerPoint演示文稿
    ppt = Presentation()
    
    # 選擇PPT版面配置（6表示空白版面）
    slide_layout = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(slide_layout)
    
    # 創建並設置標題文本框
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_box.text_frame.text = "貨架借用表"
    # 修改：使用正確的對齊方式
    title_box.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True
    
    # 計算表格的行數和列數
    rows = len(df.head(10)) + 1
    cols = len(df.columns)
    
    # 創建表格
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(5)).table
    
    # 添加表格標題行
    for col_idx, column in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(column)  # 確保轉換為字符串
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(11)
        # 修改：使用正確的對齊方式
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    
    # 添加表格數據
    for row_idx, row in enumerate(df.head(10).itertuples()):
        for col_idx, value in enumerate(row[1:]):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value) if pd.notna(value) else ""  # 處理空值
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            # 修改：使用正確的對齊方式
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    
    # 保存PowerPoint文件
    ppt.save('貨架借用表.pptx')
    print("PowerPoint 創建成功！")

except Exception as e:
    print(f"發生錯誤: {str(e)}")